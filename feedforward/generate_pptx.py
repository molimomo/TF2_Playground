from pptx import Presentation
from pptx.util import Inches
import os
result_figure_folder = './result_figures'
datasets = ['mnist','fashion_mnist']
metrices = ['val_accuracy','val_loss']
option = 'multi_runs'
measurements =['avg']
target_ranks =  [10, 20, 30]
prs = Presentation('../template_16x9.pptx')
top = Inches(1.85)
left1 = Inches(0.28)
left2 = Inches(6.67)
for dataset in datasets:
    for metric in metrices:
        for rank in target_ranks:
            title_slide_layout = prs.slide_layouts[5]
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            if metric == "val_loss":
                eval_str = "(Lower, the better)"
            else:
                eval_str = "(Higher, the better)"
            title.text = dataset + ' - rank=' + str(rank) + ', ' +str(metric) + ' ' + eval_str
            avg_config = '_'.join(['avg', metric, 'rank', str(rank)]) + '.png'
            std_config = '_'.join(['std', metric, 'rank', str(rank)]) + '.png'
            print(avg_config)
            avg_file = os.path.join(result_figure_folder,str(dataset),avg_config)
            std_file = os.path.join(result_figure_folder, str(dataset), std_config)
            pic = slide.shapes.add_picture(avg_file, left1, top)
            pic = slide.shapes.add_picture(std_file, left2, top)
prs.save('test.pptx')
