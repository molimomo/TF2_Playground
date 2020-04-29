from pptx import Presentation
from pptx.util import Inches
import os
result_figure_folder = './result_figures'
datasets = ['mnist','fashion_mnist']
kernel_sizes =[3,5,7]
metrices = ['val_accuracy','val_loss']
option = 'multi_runs'
measurements =['avg']

prs = Presentation('./template_16x9.pptx')
top = Inches(1.85)
left1 = Inches(0.28)
left2 = Inches(6.67)
for dataset in datasets:
    for metric in metrices:
        for kernel_size in kernel_sizes:
            if kernel_size == 3:
                target_ranks = [1, 2, 3]
            elif kernel_size == 5:
                target_ranks = [1, 3, 5]
            elif kernel_size == 7:
                target_ranks = [1, 3, 5, 7]
            for rank in target_ranks:
                title_slide_layout = prs.slide_layouts[5]
                slide = prs.slides.add_slide(title_slide_layout)
                title = slide.shapes.title
                if metric == "val_loss":
                    eval_str = "(Lower, the better)"
                else:
                    eval_str = "(Higher, the better)"
                title.text = dataset + ' - Kernel Size: ' + str(kernel_size) +\
                             'x'+str(kernel_size) +\
                             ', rank=' + str(rank) + ', ' +str(metric) + ' ' + eval_str
                avg_config = '_'.join(['avg', metric,
                                  'kernel_size', str(kernel_size), 'rank', str(rank)]) + '.png'
                std_config = '_'.join(['std', metric,
                                     'kernel_size', str(kernel_size), 'rank', str(rank)]) + '.png'
                print(avg_config)
                avg_file = os.path.join(result_figure_folder,str(dataset),avg_config)
                std_file = os.path.join(result_figure_folder, str(dataset), std_config)
                pic = slide.shapes.add_picture(avg_file, left1, top)
                pic = slide.shapes.add_picture(std_file, left2, top)
prs.save('test.pptx')
